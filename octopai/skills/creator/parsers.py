"""
Content parsers for different input source types.
"""

import re
import json
from pathlib import Path
from typing import Any, Optional
from abc import ABC, abstractmethod

from .models import ParsedContent


class BaseParser(ABC):
    """Base class for content parsers."""

    @abstractmethod
    def parse(self, source: dict[str, Any]) -> ParsedContent:
        """Parse content from source and return structured data."""
        pass

    def _extract_code_blocks(self, text: str) -> list[dict[str, str]]:
        """Extract code blocks from markdown text."""
        pattern = r'```(\w+)?\n(.*?)```'
        matches = re.findall(pattern, text, re.DOTALL)
        return [{"language": lang or "text", "code": code.strip()} for lang, code in matches]

    def _calculate_quality_score(self, content: str) -> float:
        """Calculate basic quality score based on content characteristics."""
        if not content:
            return 0.0

        score = 0.5

        if len(content) > 100:
            score += 0.1
        if len(content) > 500:
            score += 0.1
        if len(content) > 1000:
            score += 0.1

        has_structure = any(marker in content for marker in ['#', '##', '###', '-', '*', '1.'])
        if has_structure:
            score += 0.1

        has_code = '```' in content
        if has_code:
            score += 0.1

        return min(score, 1.0)


class TextParser(BaseParser):
    """Parser for plain text and markdown content."""

    def parse(self, source: dict[str, Any]) -> ParsedContent:
        content = source.get("content", "")
        if not content:
            return ParsedContent()

        code_blocks = self._extract_code_blocks(content)

        structured_data = {
            "headings": self._extract_headings(content),
            "paragraphs": [p.strip() for p in content.split('\n\n') if p.strip()],
            "code_blocks": code_blocks,
        }

        return ParsedContent(
            raw_text=content,
            structured_data=structured_data,
            code_blocks=code_blocks,
            metadata={"source_type": "text"},
            quality_score=self._calculate_quality_score(content),
        )

    def _extract_headings(self, text: str) -> list[dict[str, str]]:
        """Extract markdown headings with levels."""
        pattern = r'^(#{1,6})\s+(.+)$'
        headings = []
        for line in text.split('\n'):
            match = re.match(pattern, line)
            if match:
                level = len(match.group(1))
                text = match.group(2)
                headings.append({"level": level, "text": text})
        return headings


class CodeParser(BaseParser):
    """Parser for code repositories and source files."""

    def parse(self, source: dict[str, Any]) -> ParsedContent:
        file_path = source.get("file_path")
        content = source.get("content", "")

        if file_path:
            path = Path(file_path)
            if path.exists():
                content = path.read_text()
            else:
                return ParsedContent()

        if not content:
            return ParsedContent()

        code_blocks = [{"language": self._detect_language(file_path or content), "code": content}]

        structured_data = {
            "language": self._detect_language(file_path or content),
            "functions": self._extract_functions(content),
            "classes": self._extract_classes(content),
            "imports": self._extract_imports(content),
        }

        return ParsedContent(
            raw_text=content,
            structured_data=structured_data,
            code_blocks=code_blocks,
            metadata={"source_type": "code", "file_path": file_path},
            quality_score=self._calculate_quality_score(content),
        )

    def _detect_language(self, content: str) -> str:
        """Detect programming language from content or file extension."""
        if '.' in content:
            ext = content.split('.')[-1].lower()
            lang_map = {
                "py": "python",
                "js": "javascript",
                "ts": "typescript",
                "java": "java",
                "cpp": "cpp",
                "c": "c",
                "go": "go",
                "rs": "rust",
            }
            return lang_map.get(ext, "text")

        if "def " in content and "import " in content:
            return "python"
        if "function " in content and "const " in content:
            return "javascript"
        if "class " in content and "public " in content:
            return "java"

        return "text"

    def _extract_functions(self, code: str) -> list[str]:
        """Extract function names from code."""
        patterns = [
            r'def\s+(\w+)\s*\(',
            r'function\s+(\w+)\s*\(',
            r'const\s+(\w+)\s*=\s*\(',
            r'public\s+\w+\s+(\w+)\s*\(',
        ]

        functions = []
        for pattern in patterns:
            matches = re.findall(pattern, code)
            functions.extend(matches)

        return list(set(functions))

    def _extract_classes(self, code: str) -> list[str]:
        """Extract class names from code."""
        pattern = r'class\s+(\w+)'
        return re.findall(pattern, code)

    def _extract_imports(self, code: str) -> list[str]:
        """Extract import statements from code."""
        patterns = [
            r'import\s+([\w.]+)',
            r'from\s+([\w.]+)\s+import',
            r'require\([\'"]([\w.]+)[\'"]\)',
        ]

        imports = []
        for pattern in patterns:
            matches = re.findall(pattern, code)
            imports.extend(matches)

        return list(set(imports))


class DocumentParser(BaseParser):
    """Parser for document files (PDF, DOCX, etc.)."""

    def parse(self, source: dict[str, Any]) -> ParsedContent:
        file_path = source.get("file_path")
        content = source.get("content", "")

        if file_path:
            path = Path(file_path)
            if path.exists():
                ext = path.suffix.lower()
                if ext == '.pdf':
                    content = self._parse_pdf(path)
                elif ext in ['.docx', '.doc']:
                    content = self._parse_docx(path)
                else:
                    content = path.read_text()
            else:
                return ParsedContent()

        if not content:
            return ParsedContent()

        code_blocks = self._extract_code_blocks(content)

        structured_data = {
            "document_type": self._detect_document_type(file_path or ""),
            "sections": self._extract_sections(content),
            "code_blocks": code_blocks,
        }

        return ParsedContent(
            raw_text=content,
            structured_data=structured_data,
            code_blocks=code_blocks,
            metadata={"source_type": "document", "file_path": file_path},
            quality_score=self._calculate_quality_score(content),
        )

    def _parse_pdf(self, path: Path) -> str:
        """Parse PDF file (simplified implementation)."""
        try:
            import PyPDF2
            text = ""
            with open(path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except ImportError:
            return f"[PDF parsing requires PyPDF2: {path}]"

    def _parse_docx(self, path: Path) -> str:
        """Parse DOCX file (simplified implementation)."""
        try:
            from docx import Document
            doc = Document(path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            return f"[DOCX parsing requires python-docx: {path}]"

    def _detect_document_type(self, file_path: str) -> str:
        """Detect document type from file extension."""
        if file_path.endswith('.pdf'):
            return "pdf"
        elif file_path.endswith('.docx'):
            return "docx"
        elif file_path.endswith('.doc'):
            return "doc"
        return "unknown"

    def _extract_sections(self, text: str) -> list[dict[str, str]]:
        """Extract sections from document text."""
        pattern = r'^(#{1,3})\s+(.+)$'
        sections = []
        for line in text.split('\n'):
            match = re.match(pattern, line)
            if match:
                level = len(match.group(1))
                title = match.group(2)
                sections.append({"level": level, "title": title})
        return sections


class MediaParser(BaseParser):
    """Parser for audio/video content (transcription-based)."""

    def parse(self, source: dict[str, Any]) -> ParsedContent:
        file_path = source.get("file_path")
        content = source.get("content", "")
        media_type = source.get("metadata", {}).get("media_type", "unknown")

        if file_path:
            path = Path(file_path)
            if path.exists():
                content = self._transcribe_media(path, media_type)
            else:
                return ParsedContent()

        if not content:
            return ParsedContent()

        code_blocks = self._extract_code_blocks(content)

        structured_data = {
            "media_type": media_type,
            "duration_seconds": source.get("metadata", {}).get("duration", 0),
            "transcription": content,
            "segments": self._extract_segments(content),
        }

        return ParsedContent(
            raw_text=content,
            structured_data=structured_data,
            code_blocks=code_blocks,
            metadata={"source_type": "media", "file_path": file_path, "media_type": media_type},
            quality_score=self._calculate_quality_score(content),
        )

    def _transcribe_media(self, path: Path, media_type: str) -> str:
        """Transcribe audio/video content (simplified implementation)."""
        try:
            import speech_recognition as sr
            recognizer = sr.Recognizer()

            if media_type == "audio":
                with sr.AudioFile(str(path)) as source:
                    audio = recognizer.record(source)
                    return recognizer.recognize_google(audio)
            else:
                return f"[Video transcription requires additional processing: {path}]"
        except ImportError:
            return f"[Media transcription requires SpeechRecognition: {path}]"
        except Exception as e:
            return f"[Transcription error: {e}]"

    def _extract_segments(self, text: str) -> list[dict[str, str]]:
        """Extract speech segments from transcription."""
        pattern = r'\[(\d+:\d+)\]\s*(.+?)(?=\[\d+:\d+\]|$)'
        segments = []
        for match in re.finditer(pattern, text, re.DOTALL):
            timestamp = match.group(1)
            content = match.group(2).strip()
            segments.append({"timestamp": timestamp, "content": content})
        return segments


class PresentationParser(BaseParser):
    """Parser for presentation files (PPT, Keynote)."""

    def parse(self, source: dict[str, Any]) -> ParsedContent:
        file_path = source.get("file_path")
        content = source.get("content", "")

        if file_path:
            path = Path(file_path)
            if path.exists():
                ext = path.suffix.lower()
                if ext in ['.pptx', '.ppt']:
                    content = self._parse_pptx(path)
                else:
                    content = path.read_text()
            else:
                return ParsedContent()

        if not content:
            return ParsedContent()

        code_blocks = self._extract_code_blocks(content)

        structured_data = {
            "presentation_type": self._detect_presentation_type(file_path or ""),
            "slides": self._extract_slides(content),
            "code_blocks": code_blocks,
        }

        return ParsedContent(
            raw_text=content,
            structured_data=structured_data,
            code_blocks=code_blocks,
            metadata={"source_type": "presentation", "file_path": file_path},
            quality_score=self._calculate_quality_score(content),
        )

    def _parse_pptx(self, path: Path) -> str:
        """Parse PPTX file (simplified implementation)."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "---\n"
            return text
        except ImportError:
            return f"[PPTX parsing requires python-pptx: {path}]"

    def _detect_presentation_type(self, file_path: str) -> str:
        """Detect presentation type from file extension."""
        if file_path.endswith('.pptx'):
            return "pptx"
        elif file_path.endswith('.ppt'):
            return "ppt"
        return "unknown"

    def _extract_slides(self, text: str) -> list[dict[str, str]]:
        """Extract slides from presentation text."""
        slides = []
        current_slide = {"number": 1, "content": ""}

        for line in text.split('\n'):
            if line.strip() == "---":
                if current_slide["content"].strip():
                    slides.append(current_slide)
                current_slide = {"number": len(slides) + 2, "content": ""}
            else:
                current_slide["content"] += line + "\n"

        if current_slide["content"].strip():
            slides.append(current_slide)

        return slides


class TemplateParser(BaseParser):
    """Parser for skill templates."""

    def parse(self, source: dict[str, Any]) -> ParsedContent:
        file_path = source.get("file_path")
        content = source.get("content", "")

        if file_path:
            path = Path(file_path)
            if path.exists():
                content = path.read_text()
            else:
                return ParsedContent()

        if not content:
            return ParsedContent()

        code_blocks = self._extract_code_blocks(content)

        structured_data = {
            "template_name": self._extract_template_name(content),
            "variables": self._extract_variables(content),
            "sections": self._extract_template_sections(content),
            "code_blocks": code_blocks,
        }

        return ParsedContent(
            raw_text=content,
            structured_data=structured_data,
            code_blocks=code_blocks,
            metadata={"source_type": "template", "file_path": file_path},
            quality_score=self._calculate_quality_score(content),
        )

    def _extract_template_name(self, content: str) -> str:
        """Extract template name from content."""
        pattern = r'name:\s*(.+)'
        match = re.search(pattern, content)
        return match.group(1).strip() if match else ""

    def _extract_variables(self, content: str) -> list[dict[str, str]]:
        """Extract template variables."""
        pattern = r'\{\{\s*(\w+)\s*\}\}'
        variables = list(set(re.findall(pattern, content)))
        return [{"name": var, "type": "string"} for var in variables]

    def _extract_template_sections(self, content: str) -> list[dict[str, str]]:
        """Extract template sections."""
        pattern = r'^##+\s+(.+)$'
        sections = []
        for line in content.split('\n'):
            match = re.match(pattern, line)
            if match:
                sections.append({"title": match.group(1)})
        return sections
